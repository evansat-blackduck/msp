from pptx import Presentation
from pptx.util import Inches, Pt

class PPTXGenerator:
    def __init__(self, template_path):
        self.prs = Presentation(template_path)

    def add_attendees_slide(self, attendees):
        slide_layout = self.prs.slide_layouts[1]  # Assuming layout 1 is a bullet point layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Meeting Attendees"

        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame

        for attendee in attendees:
            p = tf.add_paragraph()
            p.text = attendee
            p.level = 0

    def add_success_criteria_slide(self, success_criteria):
        slide_layout = self.prs.slide_layouts[1]  # Assuming layout 1 is a bullet point layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Success Criteria"

        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame

        for criteria in success_criteria:
            p = tf.add_paragraph()
            p.text = criteria
            p.level = 0

    def save(self, output_path):
        self.prs.save(output_path)

def generate_pptx(template_path, output_path, attendees, success_criteria):
    generator = PPTXGenerator(template_path)
    generator.add_attendees_slide(attendees)
    generator.add_success_criteria_slide(success_criteria)
    generator.save(output_path)
    return output_path